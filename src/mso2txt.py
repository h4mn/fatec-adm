"""
# Microsoft Office to Text Converter

Este script Python foi desenvolvido para converter o conteúdo de arquivos do Microsoft Office em texto. Atualmente, oferece suporte aos formatos DOCX, XLSX e PPTX. O script utiliza bibliotecas como python-docx, openpyxl e python-pptx para extrair o texto dos documentos.

Autor: Hadston Nunes
Versão: 1.0
Data de Criação: 20231202

## Requisitos:
- É necessário ter o Python instalado. Neste projeto foi utilizado o Python 3.12.0;

- Foi utilizado o Virtualenv para criar um ambiente virtual para o projeto. Você pode instalá-lo usando o seguinte comando:
  pip install virtualenv

- Certifique-se de ter as bibliotecas python-docx, openpyxl e python-pptx instaladas. Você pode instalá-las usando o seguinte comando:
  pip install python-docx openpyxl python-pptx

## Como usar:
1. Execute o script fornecendo o caminho do arquivo do Microsoft Office como argumento.
   Exemplo: python office_to_text_converter.py caminho/do/arquivo.docx

2. O script extrairá o texto do arquivo e o salvará em um arquivo de texto (.txt) na mesma pasta.

Tenha em mente que este script está em desenvolvimento e pode precisar de ajustes para lidar com casos específicos. Sinta-se à vontade para contribuir ou relatar problemas no repositório do GitHub: https://github.com/h4mn/fatec-adm

## Aviso Legal:
Este script é fornecido "como está", sem garantias de qualquer tipo. Utilize-o por sua própria conta e risco.
"""

from os import path


class Caminho:
  def __init__(self, fullpath=None):
    """Inicializa os atributos do objeto
    """

    if fullpath is None:
      # O atributo path é uma referência ao caminho completo do arquivo. Ele é obtido através da função abspath da biblioteca os.
      self.fullpath = path.abspath(__file__)
    else:
      # Se o caminho fornecido não existir lança uma exceção.
      if not path.exists(fullpath):
        raise Exception("Caminho não encontrado.")

      self.fullpath = fullpath

    # O atributo dirname é uma referência ao diretório do arquivo. Ele é obtido através da função dirname da biblioteca os.
    self.dirname = path.dirname(self.fullpath)

    # O atributo parent é uma referência ao diretório pai do arquivo. Ele é obtido através da função dirname da biblioteca os.
    self.parent = path.dirname(self.dirname)

    # O atributo filename é uma referência ao nome do arquivo. Ele é obtido através da função basename da biblioteca os.
    self.filename = path.basename(self.fullpath)

    # O atributo name é uma referência ao nome do arquivo sem a extensão. Ele é obtido através da função splitext da biblioteca os.
    self.name = path.splitext(self.filename)[0]

    # O atributo ext é uma referência à extensão do arquivo. Ele é obtido através da função splitext da biblioteca os.
    self.ext = path.splitext(self.filename)[1]
    pass

  def print(self):
    """Imprime os atributos do objeto
    """
    # O atributo width é uma referência ao tamanho da maior chave do objeto. Ele é obtido através da função max da biblioteca builtins.
    width=max(len(var_name) for var_name in vars(self).keys())

    # O método print é utilizado para imprimir o caractere '=' 80 vezes.
    print('=' * 80)

    # O método vars() retorna o dicionário de atributos do objeto. O método items() retorna uma lista de tuplas contendo a chave e o valor de cada item do dicionário.
    for var_name, var_value in vars(self).items():
      # O método ljust é utilizado para alinhar o texto à esquerda, utilizando o tamanho da maior chave do objeto como referência.
      print(f"{var_name:<{width}}: {var_value}")

    print('=' * 80)
    pass


class Mso2Txt:
  def __new__(cls, mso_file, txt_file=None):
    """Converte o arquivo do Microsoft Office em um arquivo de texto
    """
    instance = super(Mso2Txt, cls).__new__(cls)

    if path.isfile(mso_file):
      instance.converter(mso_file, txt_file)
    elif path.isdir(mso_file):
      from os import listdir
      for file in listdir(mso_file):
        fullpath = path.join(mso_file, file)
        if path.exists(fullpath) and path.isfile(fullpath):
          instance.converter(fullpath, txt_file)
    pass

  def obter_txt_fullpath(self, arquivo):
    """Obtém o caminho completo do arquivo TXT
    """
    arquivo_path = Caminho(arquivo)
    txt_path = path.join(arquivo_path.dirname, "txt")
    fullpath = f"{txt_path}\{arquivo_path.name}.txt"
    return fullpath

  def obter_docx(self, arquivo):
    """Obtém o texto do arquivo DOCX
    """
    from docx import Document
    doc = Document(arquivo)
    texto = ""
    for paragraph in doc.paragraphs:
      texto += f"{paragraph.text}\n"
    return texto

  def obter_pptx(self, arquivo):
    """Obtém o texto do arquivo PPTX
    """
    from pptx import Presentation
    prs = Presentation(arquivo)
    texto = ""
    for slide in prs.slides:
      for shape in slide.shapes:
        if hasattr(shape, "text"):
          texto += f"{shape.text}\n"
    return texto

  def salvar_texto_em_arquivo(self, texto, arquivo):
    """Salva o texto em um arquivo TXT
    """
    with open(arquivo, "w", encoding="utf-8") as f:
      f.write(texto)
    pass

  def converter(self, mso_file, txt_file=None):
    """Converte o arquivo do Microsoft Office em um arquivo de texto
    """
    if txt_file is None:
      txt_file = self.obter_txt_fullpath(mso_file)
    else:
      txt_file = self.obter_txt_fullpath(txt_file)

    match path.splitext(mso_file)[1]:
      case ".docx":
        text = self.obter_docx(mso_file)
      case ".pptx":
        text = self.obter_pptx(mso_file)
      case _:
        text = None

    if text is None:
      raise Exception("Extensão não reconhecida.")

    with open(txt_file, "w", encoding="utf-8") as f:
      f.write(text)
    
    file_converted = Caminho(mso_file)
    print(f"Arquivo {file_converted.filename} convertido com sucesso.")
    pass


if __name__ == '__main__':
  """Executa o script
  """
  this = Caminho()
  pasta = path.join(this.parent, 'materiais')
  Mso2Txt(pasta)

  pass
